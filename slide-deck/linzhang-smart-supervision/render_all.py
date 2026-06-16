# -*- coding: utf-8 -*-
"""Render all 14 slides and rebuild PPTX."""
import importlib
import shutil
from pathlib import Path
from pptx import Presentation
from pptx.util import Emu

HERE = Path(__file__).parent
SLIDES = [f"slide{i:02d}" for i in range(1, 15)]

# Render
for mod in SLIDES:
    m = importlib.import_module(mod)
    tag = m.render()
    print(f"rendered {tag}")

# Build PPTX
EMU_PER_INCH = 914400
slide_w_in = 13.333
slide_h_in = 7.5
prs = Presentation()
prs.slide_width = Emu(int(slide_w_in * EMU_PER_INCH))
prs.slide_height = Emu(int(slide_h_in * EMU_PER_INCH))

png_files = sorted(HERE.glob("[0-1][0-9]-slide-*.png"))
assert len(png_files) == 14, f"expected 14 PNG files, got {len(png_files)}"

blank_layout = prs.slide_layouts[6]
for png in png_files:
    s = prs.slides.add_slide(blank_layout)
    s.shapes.add_picture(str(png), 0, 0,
                         width=prs.slide_width, height=prs.slide_height)

pptx_path = HERE / "linzhang-smart-supervision.pptx"
prs.save(pptx_path)
print(f"saved {pptx_path}")

# Copy to target location
target = Path("E:/科研/市级卫星应用专项/低空监管体系建设/临漳县低空平台/临漳县自然资源空天地一体化智能监管体系建设思路_图像式汇报版_无乱码.pptx")
target.parent.mkdir(parents=True, exist_ok=True)
shutil.copy2(pptx_path, target)
print(f"copied to {target}")
